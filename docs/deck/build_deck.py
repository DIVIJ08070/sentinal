#!/usr/bin/env python3
"""SENTINEL submission deck — Gujarat CCTV Hackathon 2026."""
import copy
from pptx import Presentation
from pptx.util import Inches as IN, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

SHOTS = "/Users/mac/Desktop/sentinel-hackathon/docs/deck/shots"
OUT = "/Users/mac/Desktop/sentinel-hackathon/deliverables/SENTINEL_Presentation.pptx"

NAVY  = RGBColor(0x0B, 0x12, 0x20)
CARD  = RGBColor(0x13, 0x1C, 0x31)
CARD2 = RGBColor(0x18, 0x23, 0x38)
LINE  = RGBColor(0x2A, 0x38, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BODY  = RGBColor(0xE6, 0xED, 0xF7)
MUT   = RGBColor(0x8F, 0xA1, 0xBC)
AMBER = RGBColor(0xF5, 0xA6, 0x23)
GREEN = RGBColor(0x34, 0xD3, 0x99)
RED   = RGBColor(0xF8, 0x71, 0x71)

FONT = "Arial"
MONO = "Courier New"
W, H = 13.333, 7.5

prs = Presentation()
prs.slide_width = IN(W)
prs.slide_height = IN(H)
BLANK = prs.slide_layouts[6]

def slide_bg(s, color=NAVY):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color

def no_line(shape):
    shape.line.fill.background()

def rect(s, x, y, w, h, fill=CARD, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06, shadow=False):
    sp = s.shapes.add_shape(shape, IN(x), IN(y), IN(w), IN(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(0.75)
    sp.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sp.adjustments[0] = radius
        except Exception: pass
    return sp

def box(s, x, y, w, h, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = s.shapes.add_textbox(IN(x), IN(y), IN(w), IN(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf

def para(tf, runs, size=12, color=BODY, bold=False, align=PP_ALIGN.LEFT,
         space_before=0, space_after=0, first=False, font=FONT, line=None):
    p = tf.paragraphs[0] if (first and not tf.paragraphs[0].runs) else tf.add_paragraph()
    p.alignment = align
    if space_before: p.space_before = Pt(space_before)
    if space_after: p.space_after = Pt(space_after)
    if line: p.line_spacing = line
    if isinstance(runs, str):
        runs = [(runs, {})]
    for text, o in runs:
        r = p.add_run(); r.text = text
        r.font.size = Pt(o.get("size", size))
        r.font.bold = o.get("bold", bold)
        r.font.italic = o.get("italic", False)
        r.font.color.rgb = o.get("color", color)
        r.font.name = o.get("font", font)
    return p

def footer(s, n):
    tf = box(s, 0.55, 7.08, 8.0, 0.3)
    para(tf, "SENTINEL — Divij Patel (Individual, Category 1) — Gujarat Police CCTV Hackathon 2026", size=9, color=MUT, first=True)
    tf2 = box(s, 11.6, 7.08, 1.18, 0.3)
    para(tf2, f"{n} / 12", size=9, color=MUT, align=PP_ALIGN.RIGHT, first=True)

def header(s, kicker, title, sub=None):
    tf = box(s, 0.55, 0.42, 12.2, 0.32)
    para(tf, kicker.upper(), size=12, color=AMBER, bold=True, first=True)
    tf = box(s, 0.55, 0.74, 12.2, 0.62)
    para(tf, title, size=25, color=WHITE, bold=True, first=True)
    if sub:
        tf = box(s, 0.55, 1.34, 12.2, 0.3)
        para(tf, sub, size=13, color=MUT, italic_hack=False, first=True) if False else para(tf, [(sub, {"italic": True})], size=13, color=MUT, first=True)

def chip(s, x, y, w, text, color=BODY, h=0.34, size=10.5, fill=CARD2, bold=False):
    c = rect(s, x, y, w, h, fill=fill, radius=0.5)
    tf = c.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = IN(0.08)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = FONT; r.font.bold = bold
    return c

def num_circle(s, x, y, n, d=0.34, fill=AMBER, color=NAVY, size=14):
    c = rect(s, x, y, d, d, fill=fill, shape=MSO_SHAPE.OVAL)
    tf = c.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(n)
    r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = color; r.font.name = FONT
    return c

def pic(s, path, x, y, w=None, h=None, border=True):
    p = s.shapes.add_picture(path, IN(x), IN(y), IN(w) if w else None, IN(h) if h else None)
    if border:
        p.line.color.rgb = LINE; p.line.width = Pt(1)
    p.shadow.inherit = False
    return p

def arrow(s, x, y, w, h=0.0, color=AMBER, weight=2.0):
    ln = s.shapes.add_connector(2, IN(x), IN(y), IN(x + w), IN(y + h))
    ln.line.color.rgb = color; ln.line.width = Pt(weight)
    le = ln.line._get_or_add_ln()
    te = le.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    le.append(te)
    return ln

# ============================================================ SLIDE 1 — TITLE
s = prs.slides.add_slide(BLANK); slide_bg(s)
tf = box(s, 0.75, 0.62, 11.8, 0.32)
para(tf, "GUJARAT POLICE CCTV HACKATHON 2026  ·  SUBMISSION DECK", size=13, color=AMBER, bold=True, first=True)
rect(s, 0.75, 1.12, 0.34, 0.34, fill=AMBER, radius=0.25)
tf = box(s, 1.28, 0.9, 11.5, 0.95)
para(tf, "SENTINEL", size=52, color=WHITE, bold=True, first=True)
tf = box(s, 0.75, 1.98, 11.8, 0.5)
para(tf, "Plate to Court in 60 Seconds — the Evidence Machine", size=22, color=AMBER, bold=True, first=True)
tf = box(s, 0.75, 2.62, 11.6, 0.85)
para(tf, [("“Where did this car go?”", {"bold": True, "color": WHITE}),
          ("  Answered from a single plate number to a court-admissible, hash-sealed evidence dossier in under a minute — live on the 30-camera government sandbox grid, on the cameras Gujarat already owns.", {})],
     size=14, color=BODY, first=True, line=1.15)
tf = box(s, 0.75, 3.52, 11.8, 0.3)
para(tf, [("PROPOSED MODEL:  ", {"bold": True, "color": MUT}),
          ("Hybrid — Model 1 (Registry & GIS)  +  Model 2/4 (Unified Viewing + Central AI Analytics & Watchlist Alerts)", {"color": BODY})],
     size=12.5, first=True)
tf = box(s, 0.75, 3.82, 11.8, 0.3)
para(tf, [("SUBMITTED BY:  ", {"bold": True, "color": MUT}),
          ("Divij Patel — Individual participant, Category 1  ·  vatsunp11@gmail.com", {"color": BODY})],
     size=12.5, first=True)
chip(s, 0.75, 4.24, 2.5, "30/30 government cameras live", color=GREEN, bold=True)
chip(s, 3.4, 4.24, 2.9, "27 plate reads on real feeds today", color=AMBER, bold=True)
chip(s, 6.45, 4.24, 2.9, "SHA-256 chain-of-custody dossier", color=BODY, bold=True)
chip(s, 9.5, 4.24, 2.4, "Zero new cameras needed", color=BODY, bold=True)
ty, tw, th = 4.82, 3.85, 2.0
for i, (img, cap) in enumerate([
        ("dashboard.png", "Live watchlist alerts — real platform, running now"),
        ("route.png", "The test-case UI — shown on the simulated-journey harness (slide 4)"),
        ("health.png", "Per-camera health — fps & bandwidth per feed")]):
    x = 0.75 + i * (tw + 0.28)
    pic(s, f"{SHOTS}/{img}", x, ty, w=tw, h=th)
    tf = box(s, x, ty + th + 0.06, tw, 0.3)
    para(tf, cap, size=9.5, color=MUT, first=True)

# ============================== SLIDE 2 — PROPOSED MODEL + KEY INNOVATIONS
s = prs.slides.add_slide(BLANK); slide_bg(s); footer(s, 2)
header(s, "Proposed model & why", "One hybrid platform, because the test case spans three models")
models = [
    ("MODEL 1 — REGISTRY & GIS  (mandatory)",
     "Every camera catalogued with GIS position, department, codec and stream URLs — synced from the gateway catalogue API. Nothing hard-coded: cameras and ids can change and the platform absorbs it."),
    ("MODEL 2 — UNIFIED VIEWING",
     "26 departments on one pane of glass: live map, HLS video wall, per-camera drawer, department filters — on existing cameras and DVR/NVRs."),
    ("MODEL 4 — CENTRAL AI + WATCHLIST ALERTS",
     "Vehicle detection + ANPR on live feeds, watchlist correlation with real-time alerts, route reconstruction, and a court-ready evidence export."),
]
my = 1.62
for title, body in models:
    rect(s, 0.55, my, 6.55, 1.06)
    tf = box(s, 0.8, my + 0.12, 6.1, 0.85)
    para(tf, title, size=13, color=AMBER, bold=True, first=True)
    para(tf, body, size=11, color=BODY, space_before=3, line=1.08)
    my += 1.2
rect(s, 0.55, my + 0.02, 6.55, 1.34, fill=CARD2)
tf = box(s, 0.8, my + 0.14, 6.1, 1.14)
para(tf, "WHY HYBRID", size=12, color=WHITE, bold=True, first=True)
para(tf, "The evaluation exercises all three at once: onboard ~50 heterogeneous cameras (M1), monitor them centrally (M2), then trace a designated plate and alert on a watchlist (M4). One registry-driven platform covers the full graded path with no duplicated infrastructure.",
     size=11, color=BODY, space_before=3, line=1.1)
tf = box(s, 7.4, 1.56, 5.35, 0.3)
para(tf, "KEY INNOVATIONS", size=13, color=WHITE, bold=True, first=True)
innov = [
    ("PTS-anchored evidence clock", "every timestamp derives from stream PTS, never wall-clock or declared fps — survives the gateway's own timing traps."),
    ("Confusion-tolerant plate matcher", "0↔O, 1↔I, 5↔S, 8↔B misreads recovered at displayed confidence — never silently merged."),
    ("Physics plausibility filter", "impossible hops rejected with plain-language reasons (“214 km/h over 3.1 km — discarded”)."),
    ("Hash-chained evidence dossier", "SHA-256 chain-of-custody PDF: snapshots, GPS/timestamp table, operator identity, audit trail."),
    ("Edge-first 80,000-camera design", "video stays on departmental NVRs; ~1–3 Kbps per camera of metadata travels upstream."),
]
iy = 1.94
for i, (t, b) in enumerate(innov, 1):
    num_circle(s, 7.4, iy + 0.02, i, d=0.3, size=12)
    tf = box(s, 7.84, iy, 4.95, 0.95)
    para(tf, [(t + " — ", {"bold": True, "color": WHITE}), (b, {"color": BODY})], size=10.5, first=True, line=1.05)
    iy += 1.0

# ==================================== SLIDE 3 — THE 60-SECOND STORY (workflow)
s = prs.slides.add_slide(BLANK); slide_bg(s); footer(s, 3)
header(s, "Solution overview · end-to-end workflow", "The 60-second story — five beats, one arc, the exact live demo")
beats = [
    ("COLD OPEN", "All 30 government feeds live on one pane of glass. Health board green, per-feed bandwidth counters visible. No login screen, no loading."),
    ("THE ASK", "Judge reads a plate. Typed once. The route animates camera-by-camera: timestamps, plate-crop evidence cards, per-hop confidence chips."),
    ("PROOF OF LIFE", "The same plate crosses a live camera — the watchlist alert fires with an audible ping; the operator's acknowledgment lands in the audit log on screen."),
    ("CHAOS BEAT", "We kill a feed on purpose. Amber → reconnecting → green — the same backoff code that survives the sandbox's looping feeds."),
    ("PLATE TO COURT", "One click: chain-of-custody PDF — hashed frames, GPS/timestamp table, operator identity. Handed to the panel, printed."),
]
bw, bh, gap = 2.32, 3.5, 0.13
bx, by = 0.55, 1.75
for i, (t, b) in enumerate(beats, 1):
    x = bx + (i - 1) * (bw + gap)
    rect(s, x, by, bw, bh)
    num_circle(s, x + 0.16, by + 0.16, i, d=0.38, size=15)
    tf = box(s, x + 0.16, by + 0.68, bw - 0.32, bh - 0.85)
    para(tf, t, size=13, color=AMBER, bold=True, first=True)
    para(tf, b, size=10.5, color=BODY, space_before=5, line=1.12)
rect(s, 0.55, 5.6, 12.23, 0.78, fill=CARD2)
tf = box(s, 0.85, 5.74, 11.7, 0.55)
para(tf, [("Rehearsal protocol: ", {"bold": True, "color": WHITE}),
          ("every beat rehearsed cold-start daily; three timed dress runs (real grid / injected failure / fully offline). Beat 2 is never canned — the judge-plate trace is always live.", {})],
     size=11.5, color=BODY, first=True, line=1.1)

# ============================ SLIDE 4 — CRITERION 1: GOVERNMENT TEST CASE
s = prs.slides.add_slide(BLANK); slide_bg(s); footer(s, 4)
header(s, "Criterion 1 — the government test case", "The plate trace, proven on the real grid — not hoped for")
pic(s, f"{SHOTS}/route.png", 0.55, 1.62, w=7.55, h=4.13)
tf = box(s, 0.55, 5.82, 7.55, 0.35)
para(tf, "The route UI on the platform's simulated-journey demo harness (snapshots watermarked SIMULATED FEED) — the same engine serves real-grid queries: e.g. 5 timestamped sightings of CMCI801 on cam23 with 4 OCR misreads fuzzy-recovered.",
     size=10, color=MUT, first=True, line=1.1)
facts = [
    ("ONBOARDED LIVE — SEPT 2", "The full government sandbox grid via catalogue sync: 30 heterogeneous cameras — 24 H.264 + 6 HEVC, five resolutions, declared 10–30 fps. RTSP forced over TCP, max 4 concurrent captures (pacing rule).", BODY),
    ("MEASURED, NOT TRUSTED", "cam01 declares 30 fps but delivers ~12.5 fps (median PTS gap exactly 80 ms) — the integration guide's own warning, demonstrated live. All timing is PTS-anchored; never arrival time, never declared fps.", AMBER),
    ("REAL DETECTIONS TODAY", "293 detections on 11 real cameras, incl. 27 plate reads on 7 of them — each with a PTS-anchored UTC timestamp, snapshot and bbox, stored where GET /api/detections serves them.", GREEN),
]
fy = 1.62
for t, b, tc in facts:
    rect(s, 8.3, fy, 4.48, 1.5)
    tf = box(s, 8.52, fy + 0.13, 4.05, 1.3)
    para(tf, t, size=12, color=tc, bold=True, first=True)
    para(tf, b, size=10, color=BODY, space_before=3, line=1.08)
    fy += 1.63
rect(s, 8.3, fy - 0.06, 4.48, 0.62, fill=CARD2)
tf = box(s, 8.52, fy + 0.03, 4.05, 0.5)
para(tf, [("Score-sheet line: ", {"bold": True, "color": WHITE}),
          ("“identify, trace and present the complete timestamped route” → demo beat 2.", {})],
     size=10, color=BODY, first=True, line=1.05)

# ================================ SLIDE 5 — CRITERION 5: AI ANALYTICS
s = prs.slides.add_slide(BLANK); slide_bg(s); footer(s, 5)
header(s, "Criterion 5 — AI analytics approach", "We show you what the system refused to believe — and why")
stages = ["RTSP over TCP\nPTS-anchored frames", "YOLOv8n\nvehicle detection", "Plate localization\nYOLO-v9-t 384", "fast-plate-ocr\nplate reading", "Normalize + match\nexact / fuzzy, ranked"]
sw, sh, sgap = 2.24, 0.86, 0.26
sx, sy = 0.55, 1.66
for i, st in enumerate(stages):
    x = sx + i * (sw + sgap)
    rect(s, x, sy, sw, sh, fill=CARD2)
    tf = box(s, x + 0.1, sy, sw - 0.2, sh, anchor=MSO_ANCHOR.MIDDLE)
    lines = st.split("\n")
    para(tf, lines[0], size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
    para(tf, lines[1], size=9.5, color=MUT, align=PP_ALIGN.CENTER)
    if i < 4:
        arrow(s, x + sw + 0.02, sy + sh / 2, 0.22)
tf = box(s, 0.55, 2.62, 12.2, 0.3)
para(tf, "CPU-only inference, per-camera buffers (no fixed-shape batch across heterogeneous cameras), detector reset on loop discontinuities.",
     size=10.5, color=MUT, first=True)
rect(s, 0.55, 3.1, 6.0, 2.6)
tf = box(s, 0.85, 3.28, 5.45, 2.3)
para(tf, "THE 0.93 RECOVERY", size=14, color=GREEN, bold=True, first=True)
para(tf, [("OCR reads ", {}), ("GJ01A81234", {"font": MONO, "bold": True, "color": WHITE}),
          (" → matched ", {}), ("GJ01AB1234", {"font": MONO, "bold": True, "color": GREEN}),
          (" at 0.93, flagged fuzzy.", {})], size=11.5, color=BODY, space_before=7, line=1.15)
para(tf, "8 misread as B — an exact-match system shows a hole in the route; SENTINEL shows a recovered sighting at displayed confidence. Ranked, flagged, never silently merged; an all-fuzzy route says so in a banner.",
     size=11, color=BODY, space_before=7, line=1.15)
rect(s, 6.78, 3.1, 6.0, 2.6)
tf = box(s, 7.08, 3.28, 5.45, 2.3)
para(tf, "THE REJECTED HOP", size=14, color=RED, bold=True, first=True)
para(tf, [("“Implied speed 214 km/h over 3.1 km — physically impossible, discarded as false ANPR match.”", {"italic": True, "color": WHITE})],
     size=11.5, space_before=7, line=1.15)
para(tf, "A physics filter over GIS distance / PTS-anchored Δt greys the hop on the map with a plain-language reason; a retro-rejection guard stops a false first sighting from poisoning the route.",
     size=11, color=BODY, space_before=7, line=1.15)
rect(s, 0.55, 5.9, 12.23, 0.72, fill=CARD2)
tf = box(s, 0.85, 6.02, 11.7, 0.52)
para(tf, [("Uncounterfeitable: ", {"bold": True, "color": AMBER}),
          ("this proves the system models vehicles moving through the world — not strings matching in a database. Face recognition: roadmap only, deliberately not demoed.", {})],
     size=11.5, color=BODY, first=True, line=1.1)

# ============ SLIDE 6 — WATCHLIST CORRELATION + REAL-TIME ALERTS
s = prs.slides.add_slide(BLANK); slide_bg(s); footer(s, 6)
header(s, "Watchlist correlation & real-time alerts", "Every detection, against every active entry — in real time")
steps = [
    ("NORMALIZE", "One shared rule for every plate everywhere: uppercase, A–Z and 0–9 only. The watchlist and the detector can never disagree on format."),
    ("CORRELATE", "Each detection is matched against all active watchlist entries: exact on the normalized plate, plus fuzzy via Indian-plate canonicalization + weighted edit distance over OCR-confusion pairs (0↔O, 1↔I, 5↔S, 8↔B, 6↔G, 2↔Z cost 0.25; any other edit 1.0; total ≤ 1.0 matches) — always labelled fuzzy with displayed confidence (a single confusion misread scores 0.93)."),
    ("ALERT IN UNDER A SECOND", "Matches push over WebSocket to every operator screen: audible ping, snapshot evidence card, category and priority badges, map pan to the camera."),
    ("ACKNOWLEDGE & AUDIT", "The operator's acknowledgment lands in the append-only audit log — the same log the evidence dossier cites. Nothing disappears silently."),
]
sy = 1.66
for i, (t, b) in enumerate(steps, 1):
    rect(s, 0.55, sy, 8.6, 1.18)
    num_circle(s, 0.75, sy + 0.16, i, d=0.34, size=13)
    tf = box(s, 1.28, sy + 0.13, 7.7, 0.98)
    para(tf, t, size=12.5, color=AMBER, bold=True, first=True)
    para(tf, b, size=10.5, color=BODY, space_before=2, line=1.08)
    sy += 1.3
pic(s, f"{SHOTS}/alerts_crop.png", 9.4, 1.66, w=3.35, h=5.0)
tf = box(s, 9.4, 6.72, 3.35, 0.3)
para(tf, "Live alert cards on the running platform.", size=9.5, color=MUT, first=True)

# ======================= SLIDE 7 — CRITERION 3: ARCHITECTURE + TECH
s = prs.slides.add_slide(BLANK); slide_bg(s); footer(s, 7)
header(s, "Criterion 3 — high-level design", "Tiered edge architecture: video stays where it lives")
tiers = [
    (0.55, "80,000 CAMERAS + DVR/NVRs", "26 departments' existing CCTV estate — untouched. Video is decoded and analysed near the camera and never leaves departmental storage.", "video stays local"),
    (5.0, "EDGE / REGIONAL PoPs (~30)", "Decode + ANPR run locally per district PoP; per-camera health metrics measured at the edge; store-and-forward for low-connectivity sites.", "~1–3 Kbps/cam metadata →"),
    (9.45, "STATE CENTRE", "Registry + GIS · watchlist correlation · route engine · evidence dossier · append-only audit. Metadata-scale only — no video backhaul.", ""),
]
for x, t, b, lab in tiers:
    rect(s, x, 1.7, 3.35, 2.5)
    tf = box(s, x + 0.22, 1.88, 2.92, 2.2)
    para(tf, t, size=13, color=WHITE, bold=True, first=True)
    para(tf, b, size=10.5, color=BODY, space_before=6, line=1.12)
arrow(s, 3.95, 2.95, 1.0); arrow(s, 8.4, 2.95, 1.0)
tf = box(s, 3.7, 2.5, 1.55, 0.35)
para(tf, "metadata", size=9.5, color=AMBER, align=PP_ALIGN.CENTER, first=True)
tf = box(s, 8.15, 2.5, 1.55, 0.35)
para(tf, "~1–3 Kbps/cam", size=9.5, color=AMBER, align=PP_ALIGN.CENTER, first=True)
rect(s, 0.55, 4.4, 12.23, 0.92, fill=CARD2)
tf = box(s, 0.85, 4.52, 11.7, 0.72)
para(tf, [("Interoperability: ", {"bold": True, "color": WHITE}),
          ("the catalogue registry API is the contract — cameras and ids can change and the platform absorbs it; per-camera codec/resolution read from the catalogue, no uniform-grid assumption. VAHAN / SARTHI / eGujCop integration readiness documented (HLD §10). Evolution path on the same contracts: SQLite → PostgreSQL+PostGIS, REST → Kafka, threads → K8s per PoP.", {})],
     size=10.5, color=BODY, first=True, line=1.12)
tf = box(s, 0.55, 5.52, 12.2, 0.3)
para(tf, "TECHNOLOGIES USED", size=11, color=MUT, bold=True, first=True)
tech = ["Python 3.12", "FastAPI", "SQLAlchemy", "OpenCV + FFmpeg (RTSP/TCP)", "Ultralytics YOLOv8n", "fast-plate-ocr", "open-image-models", "React 18 + Vite", "Leaflet GIS", "hls.js", "WebSocket", "JWT RBAC", "SQLite → PostGIS"]
tx, ty2 = 0.55, 5.88
for t in tech:
    wch = 0.32 + len(t) * 0.082
    if tx + wch > 12.8:
        tx = 0.55; ty2 += 0.44
    chip(s, tx, ty2, wch, t, color=BODY, size=10)
    tx += wch + 0.16

# ======================= SLIDE 8 — CRITERION 6: THE 80K ARITHMETIC
s = prs.slides.add_slide(BLANK); slide_bg(s); footer(s, 8)
header(s, "Criterion 6 — scalability to 80,000 cameras", "Computed, not asserted — and measured live on the health board")
rect(s, 0.55, 1.7, 6.0, 3.3)
tf = box(s, 0.85, 1.92, 5.45, 3.0)
para(tf, "CENTRALIZED (THE WRONG WAY)", size=13, color=RED, bold=True, first=True)
para(tf, "160 Gbps", size=54, color=RED, bold=True, space_before=10)
para(tf, "80,000 cameras × 2 Mbps of video hauled into one site. Unbuildable on GSWAN — and pointless to build.",
     size=12, color=BODY, space_before=10, line=1.15)
rect(s, 6.78, 1.7, 6.0, 3.3)
tf = box(s, 7.08, 1.92, 5.45, 3.0)
para(tf, "EDGE-FIRST (SENTINEL)", size=13, color=GREEN, bold=True, first=True)
para(tf, "80–240 Mbps", size=54, color=GREEN, bold=True, space_before=10)
para(tf, "Detection metadata only: ~1–3 Kbps per camera upstream — a 650× to 2,000× reduction. Video stays on departmental NVRs, exactly as today.",
     size=12, color=BODY, space_before=10, line=1.15)
rect(s, 0.55, 5.28, 12.23, 1.3, fill=CARD2)
tf = box(s, 0.85, 5.44, 11.7, 1.05)
para(tf, [("Shown live, on screen: ", {"bold": True, "color": AMBER}),
          ("the camera-health board puts both numbers side by side for the real grid — per-camera video bandwidth staying at the edge next to the metadata actually travelling upstream (rows are measured in real time for every camera with a live worker attached).", {})],
     size=12, color=BODY, first=True, line=1.15)
para(tf, "Honest scope: architecture plus measured ratios on one machine — not a deployed fleet. GPU sizing, analytics tiering and retention are stated as explicit assumptions in the HLD (§6).",
     size=11, color=MUT, space_before=6, line=1.12)

# ======================= SLIDE 9 — CRITERION 4: PLATFORM MATURITY
s = prs.slides.add_slide(BLANK); slide_bg(s); footer(s, 9)
header(s, "Criterion 4 — platform maturity & security", "Built for the 2 a.m. shift — resilience shown, not claimed")
cards = [
    ("HEALTH BOARD", "Per-feed measured fps, last-frame age, reconnect count and live bandwidth — problem feeds sort first; auto-alert on feed drop."),
    ("KILL-A-FEED, LIVE", "We cut a feed on stage: green → amber → reconnecting → green. Exponential backoff 2 s → 30 s plus loop-discontinuity reset — rehearsed daily."),
    ("RBAC + AUDIT", "Token-authenticated viewer / operator / admin roles gate watchlist changes, acks and dossier export; append-only audit log; the identity on the dossier comes from the token, never a spoofable header."),
    ("TRIPLE-LAYER FALLBACK", "DEMO_MODE flag, a fully local 50-camera gateway needing zero venue network, and one dress rehearsal run entirely offline."),
]
cy = 1.66
for t, b in cards:
    rect(s, 0.55, cy, 9.1, 1.18)
    tf = box(s, 0.82, cy + 0.13, 8.6, 0.98)
    para(tf, t, size=12.5, color=AMBER, bold=True, first=True)
    para(tf, b, size=10.5, color=BODY, space_before=2, line=1.08)
    cy += 1.3
pic(s, f"{SHOTS}/health_crop.png", 9.9, 1.66, w=2.42, h=5.26)
tf = box(s, 9.9, 6.94, 2.6, 0.26)
para(tf, "Health board: 30/30 streams reporting.", size=9, color=MUT, first=True)

# ======================= SLIDE 10 — COST + OPERATIONAL BENEFITS
s = prs.slides.add_slide(BLANK); slide_bg(s); footer(s, 10)
header(s, "Cost & operational benefits for policing", "You do not need a single new camera")
rect(s, 0.55, 1.7, 5.6, 3.15)
tf = box(s, 0.85, 1.95, 5.05, 2.8)
para(tf, "₹3–7k", size=58, color=AMBER, bold=True, first=True)
para(tf, "per EXISTING camera to bring it onto SENTINEL", size=13, color=WHITE, bold=True, space_before=4)
para(tf, [("vs ", {}), ("₹25k+", {"bold": True, "color": RED}), (" per new camera installed.", {})],
     size=13, color=BODY, space_before=10)
para(tf, "Video never leaves departmental storage — no new backhaul, no forklift replacement.",
     size=11.5, color=MUT, space_before=8, line=1.12)
bens = [
    ("8–16 hours → sub-second", "a vehicle trace that today costs an operator a shift of manual DVR review becomes one query."),
    ("The watchlist never sleeps", "every live feed cross-referenced 24×7; alerts reach the control room in under a second, with evidence attached."),
    ("Evidence that survives court", "the dossier reaches the charge sheet hash-sealed, with operator identity and audit trail — CCTV evidence stops dying between control room and court."),
    ("Phased, low-risk rollout", "P1 pilot: one commissionerate, ~2,000 cameras, ~₹1.5 cr — validates every ratio before scaling. Statewide sketch ₹15–20 cr excl. cameras (HLD §11, ±30%)."),
]
by = 1.7
for t, b in bens:
    rect(s, 6.38, by, 6.4, 1.02)
    tf = box(s, 6.64, by + 0.11, 5.9, 0.85)
    para(tf, [(t + " — ", {"bold": True, "color": WHITE}), (b, {"color": BODY})], size=11, first=True, line=1.1)
    by += 1.14
rect(s, 0.55, 5.05, 5.6, 1.1, fill=CARD2)
tf = box(s, 0.85, 5.2, 5.05, 0.85)
para(tf, [("Same infrastructure amortizes every future analytic ", {"bold": True, "color": WHITE}),
          ("at software-only cost — the brief's “existing infrastructure to the maximum practical extent”, done literally.", {})],
     size=11, color=BODY, first=True, line=1.12)

# ======================= SLIDE 11 — SCORE-SHEET MIRROR (TICK MATRIX)
s = prs.slides.add_slide(BLANK); slide_bg(s); footer(s, 11)
header(s, "The score sheet, mirrored", "Seven criteria + every bonus line — each tick names its proof")
crit = [
    ("1 · Government test case", "designated plate → timestamped route, live on the real grid (slide 4, demo beat 2)"),
    ("2 · Presentation", "this deck: one slide per criterion, transcribable straight into the rubric"),
    ("3 · High-level design", "tiered edge architecture; 12-section HLD submitted as PDF (slide 7)"),
    ("4 · Platform maturity", "health board + on-stage kill-a-feed recovery + RBAC/audit (slide 9)"),
    ("5 · Analytics quality", "physics-filter rejection + 0.93 fuzzy recovery, confidence shown (slide 5)"),
    ("6 · Scalability to 80k", "160 Gbps vs 80–240 Mbps — computed and measured (slide 8)"),
    ("7 · Completeness", "deck, HLD, two videos, timestamped report, repo — all banked early (slide 12)"),
]
bonus = [
    ("Multi-camera correlation", "route engine w/ physics filter + retro-rejection guard (regression-tested); fuzzy recoveries proven on the real grid"),
    ("Cybersecurity / audit / RBAC", "SHA-256 dossier chain, JWT roles, append-only audit the export cites"),
    ("Edge / bandwidth efficiency", "~1–3 Kbps/cam upstream; ratio measured live on the health board"),
    ("Camera-health monitoring", "per-feed fps / frame-age / reconnects / bandwidth + drop alerts"),
    ("Existing-infrastructure reuse", "zero new cameras; DVR/NVRs untouched; ₹3–7k per existing camera"),
    ("Real-time alerting", "watchlist alert < 1 s over WebSocket, snapshot attached, ack audited"),
]
tf = box(s, 0.55, 1.6, 6.0, 0.3)
para(tf, "THE SEVEN EVALUATION CRITERIA", size=12, color=MUT, bold=True, first=True)
ry = 1.96
for t, b in crit:
    rect(s, 0.55, ry, 6.1, 0.62, fill=CARD)
    tf = box(s, 0.72, ry + 0.06, 0.35, 0.5)
    para(tf, "✓", size=15, color=GREEN, bold=True, first=True)
    tf = box(s, 1.12, ry + 0.05, 5.4, 0.55)
    para(tf, t, size=10.5, color=WHITE, bold=True, first=True)
    para(tf, b, size=8.5, color=MUT, line=1.0)
    ry += 0.71
tf = box(s, 6.95, 1.6, 5.8, 0.3)
para(tf, "BONUS CRITERIA — NO UNPROVEN TICKS", size=12, color=MUT, bold=True, first=True)
ry = 1.96
for t, b in bonus:
    rect(s, 6.95, ry, 5.85, 0.62, fill=CARD)
    tf = box(s, 7.12, ry + 0.06, 0.35, 0.5)
    para(tf, "✓", size=15, color=GREEN, bold=True, first=True)
    tf = box(s, 7.52, ry + 0.05, 5.15, 0.55)
    para(tf, t, size=10.5, color=WHITE, bold=True, first=True)
    para(tf, b, size=8.5, color=MUT, line=1.0)
    ry += 0.71
rect(s, 6.95, ry, 5.85, 0.62, fill=CARD2)
tf = box(s, 7.12, ry + 0.07, 5.5, 0.5)
para(tf, [("Not claimed: ", {"bold": True, "color": AMBER}),
          ("face recognition is roadmap, not demo; every legal statement scoped as design provision.", {})],
     size=9.5, color=BODY, first=True, line=1.05)

# ============================================= SLIDE 12 — CLOSE
s = prs.slides.add_slide(BLANK); slide_bg(s); footer(s, 12)
tf = box(s, 1.2, 2.0, 10.9, 1.9, anchor=MSO_ANCHOR.MIDDLE)
para(tf, "“Plate to court in sixty seconds —", size=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
para(tf, [("on the cameras Gujarat already owns.”", {"color": AMBER})], size=36, bold=True, align=PP_ALIGN.CENTER, space_before=6)
tf = box(s, 1.2, 4.05, 10.9, 0.35)
para(tf, "EVERYTHING THE CHECKLIST ASKS FOR — BANKED EARLY, SUBMITTED BEFORE NOON ON 7 SEPTEMBER", size=11, color=MUT, bold=True, align=PP_ALIGN.CENTER, first=True)
arts = [
    ("DECK", "this file + PDF"),
    ("HLD", "12-section PDF"),
    ("VIDEO ×2", "own-feed + government-feed"),
    ("REPORT", "timestamped = dossier export"),
    ("REPO", "clean-clone quickstart, make test"),
    ("LIVE DEMO", "cold-start rehearsed, offline fallback"),
]
aw, agap = 1.92, 0.12
ax = (W - (aw * 6 + agap * 5)) / 2
for i, (t, b) in enumerate(arts):
    x = ax + i * (aw + agap)
    rect(s, x, 4.55, aw, 1.15)
    tf = box(s, x + 0.12, 4.55, aw - 0.24, 1.15, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, t, size=12.5, color=AMBER, bold=True, align=PP_ALIGN.CENTER, first=True)
    para(tf, b, size=9, color=BODY, align=PP_ALIGN.CENTER, space_before=3, line=1.0)
tf = box(s, 1.2, 6.1, 10.9, 0.4)
para(tf, "SENTINEL  ·  Hybrid Model 1 + 2/4  ·  Divij Patel (Individual, Category 1)  ·  Gujarat Police CCTV Hackathon 2026", size=12, color=MUT, align=PP_ALIGN.CENTER, first=True)

prs.core_properties.title = "SENTINEL — Gujarat Police CCTV Hackathon 2026 Submission Deck"
prs.core_properties.author = "Divij Patel"
prs.core_properties.subject = "Gujarat Police CCTV Hackathon 2026 — Category 1, Individual"
prs.save(OUT)
print("saved", OUT, "slides:", len(prs.slides.__iter__.__self__._sldIdLst))
